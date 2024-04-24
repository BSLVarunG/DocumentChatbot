#imports
#for AICore connectivity and triggering LLM
from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate
from gen_ai_hub.proxy.langchain.init_models import init_llm
from gen_ai_hub.proxy.langchain.init_models import init_embedding_model

#for user inputs
import io

#to convert PDF files to necessary format
import PyPDF2

#to convert Word Documents to necessary format
from docx import Document
from prompt_toolkit.shortcuts import PromptSession
from prompt_toolkit.lexers import PygmentsLexer
from pygments.lexers import PythonLexer

#to convert PPT to necessary format
from pptx import Presentation

#to convert Email to necessary format
import extract_msg
import datetime

file_dict = {
'word':'/home/user/projects/TestCode/TestFiles/TestFiles/EssayOnPython.docx',
'pdf':'/home/user/projects/TestCode/TestFiles/TestFiles/Essay On AI.pdf',
'ppt':'/home/user/projects/TestCode/TestFiles/TestFiles/TimeManagement1.pptx',
'email':'/home/user/projects/TestCode/TestFiles/TestFiles/Test Email.msg'
}

def word2textconv(file_path):
    document_text = ''
    doc = Document(file_path)
    for paragraph in doc.paragraphs:
        document_text += paragraph.text + "\n"
    return(document_text)

def pdf2textconv(file_path):
    document_text = ''
    with open(file_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        num_pages = len(pdf_reader.pages)
        for page_num in range(num_pages):
            page = pdf_reader.pages[page_num]
            document_text += page.extract_text()
    return document_text

def ppt2textconv(file_path):
    document_text = ''
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                document_text += shape.text + '\n'
    return document_text

def email2textconv(file_path):
    document_text = ''
    msg = extract_msg.Message(file_path)
    document_text += msg.sender+'\n'
    document_text += (msg.date).strftime("%Y-%m-%d %H:%M:%S")+'\n'
    document_text += msg.subject+'\n'
    document_text += msg.body+'\n'
    msg.close()
    return document_text

def LLMSummarize(document_text):
    summarize_prompt = "Your task is to summarize the following document in a single sentence of max length of 15 words\n"+document_text
    llm = init_llm('gpt-35-turbo', temperature=0., max_tokens=3000) 
    summary = llm.invoke(summarize_prompt).content
    return summary 

def LLMQnA(file_text,question):
    template = 'You are a Chatbot assisstant. Your task is to refer to the document below and answer questions based on the document\n'+file_text+'\nQuestion: Can You answer questions based on the document above?\nAnswer: Yes\nQuestion: {question}\nAnswer:' 
    Prompt = PromptTemplate(template=template, input_variables=['question'])
    llm = init_llm('gpt-35-turbo', temperature=0., max_tokens=3000) 
    llm_chain = LLMChain(prompt=Prompt, llm=llm) 
    response = llm_chain.invoke(question) 
    return response['text']

def main():
    intro = 'Hello Im the Birlasoft Text Summarization module. I can help you with summarization and answering questions of 4 types of documents namely: Word documents, Presentations, PDF files and Email attachments\n'
    print(intro)
    menu = 'Enter 1 for Word Documents\nEnter 2 for PDF files\nEnter 3 for Presentations\nEnter 4 for Emails\nEnter 5 to exit the program\n'
    choice = input(menu)
    while True:
        with open("Chathistory.txt",'w') as file:
            file.write(intro+menu+choice)
            if choice == '1':
                file_text = word2textconv(file_dict['word'])
                file.write('Reading the word document... ')
                print('Reading the word document... ')
                document_summary = LLMSummarize(file_text)
                print('The summary of the document is as follows'+document_summary)
                file.write('The summary of the document is as follows'+document_summary +'\n')
                question = input('What would you like to know about this document?')
                file.write('Question:'+question+'\n')
                while True:
                    query_response = LLMQnA(file_text,question)
                    file.write('Answer:'+query_response +'\n')
                    print(query_response)
                    question=input('Is there anything else that you need to know about the document?')
                    if(question.lower()=='no'):
                        break
                if(question.lower()=='no'):
                    print('Goodbye!')
                    file.write('Goodbye\n') 
                    break
                choice =  input(menu)   
            elif choice == '2':
                file_text = pdf2textconv(file_dict['pdf'])
                file.write("Reading the PDF document... ")
                print("Reading the PDF document... ")
                document_summary = LLMSummarize(file_text)
                print('The summary of the document is as follows'+document_summary)
                file.write('The summary of the document is as follows'+document_summary +'\n')
                question = input('What would you like to know about this document?')
                file.write(question+'\n')
                while True:
                    query_response = LLMQnA(file_text,question)
                    file.write('Answer:'+query_response +'\n')
                    print(query_response)
                    question=input('Is there anything else that you need to know about the document?')
                    if(question.lower()=='no'):
                        break
                if(question.lower()=='no'):
                    print('Goodbye!')
                    file.write('Goodbye\n') 
                    break
                choice =  input(menu)
            elif choice == '3':
                file_text = ppt2textconv(file_dict['ppt'])
                file.write("Reading the PPT document... ")
                print("Reading the PPT document... ")
                document_summary = LLMSummarize(file_text)
                print('The summary of the document is as follows'+document_summary)
                file.write('The summary of the document is as follows'+document_summary +'\n')
                question = input('What would you like to know about this document?')
                file.write(question+'\n')
                while True:
                    query_response = LLMQnA(file_text,question)
                    file.write('Answer:'+query_response +'\n')
                    print(query_response)
                    question=input('Is there anything else that you need to know about the document?')
                    if(question.lower()=='no'):
                        break
                if(question.lower()=='no'):
                    print('Goodbye!')
                    file.write('Goodbye\n') 
                    break
                choice =  input(menu)
            elif choice == '4':
                file_text = email2textconv(file_dict['email'])
                file.write("Reading the Email attachment... ")
                document_summary = LLMSummarize(file_text)
                print('The summary of the document is as follows'+document_summary)
                file.write('The summary of the document is as follows'+document_summary +'\n')
                question = input('What would you like to know about this document?')
                file.write(question+'\n')
                while True:
                    query_response = LLMQnA(file_text,question)
                    file.write('Answer:'+query_response +'\n')
                    print(query_response)
                    question=input('Is there anything else that you need to know about the document?')
                    if(question.lower()=='no'):
                        break
                if(question.lower()=='no'):
                    print('Goodbye!')
                    file.write('Goodbye\n') 
                    break
                choice =  input(menu)
            elif choice=='5':
                print('Goodbye!')
                file.write('Goodbye\n') 
                break   
            else:
                print("Invalid input")
                choice = input(menu)
        file.write('Chat Terminated\n')
    print('Program Terminated')

if __name__ == "__main__":
    main()
